#!/usr/bin/env python3
"""Generate the two-slide Cursor workflow training deck."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs/training/cursor-workflow.pptx"

WIDE_WIDTH = Inches(13.333)
WIDE_HEIGHT = Inches(7.5)

BG = RGBColor(0x07, 0x11, 0x1F)
SURFACE = RGBColor(0x0F, 0x1C, 0x2E)
SURFACE_2 = RGBColor(0x12, 0x23, 0x38)
BORDER = RGBColor(0x26, 0x39, 0x4D)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x9A, 0xA9, 0xBC)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
CYAN = RGBColor(0x38, 0xBD, 0xF8)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0xA7, 0x8B, 0xFA)

FONT = "Inter"
MONO = "DejaVu Sans Mono"


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=WHITE,
    bold=False,
    font=FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.text = text
    paragraph.font.name = font
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return shape


def add_rect(
    slide,
    x,
    y,
    w,
    h,
    *,
    fill=SURFACE,
    line=BORDER,
    radius=True,
    line_width=1,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        kind, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def add_chevron(slide, x, y, color=MUTED):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.18), Inches(0.34)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, number, title, subtitle, x, width, accent):
    add_rect(slide, x, 1.72, width, 1.48)
    bar = add_rect(
        slide,
        x,
        1.72,
        width,
        0.07,
        fill=accent,
        line=accent,
        radius=False,
        line_width=0,
    )
    bar.line.fill.background()
    add_text(
        slide,
        number,
        x + 0.15,
        1.88,
        0.48,
        0.34,
        size=12,
        color=accent,
        bold=True,
        font=MONO,
    )
    add_text(
        slide,
        title,
        x + 0.15,
        2.16,
        width - 0.3,
        0.44,
        size=16,
        bold=True,
        valign=MSO_ANCHOR.TOP,
    )
    add_text(
        slide,
        subtitle,
        x + 0.15,
        2.63,
        width - 0.3,
        0.39,
        size=10.5,
        color=MUTED,
        valign=MSO_ANCHOR.TOP,
    )


def add_header(slide, section, title, subtitle, slide_number):
    add_rect(
        slide,
        0,
        0,
        13.333,
        7.5,
        fill=BG,
        line=BG,
        radius=False,
        line_width=0,
    )
    add_rect(
        slide,
        0,
        0,
        13.333,
        0.06,
        fill=GREEN,
        line=GREEN,
        radius=False,
        line_width=0,
    )
    add_text(
        slide,
        section.upper(),
        0.52,
        0.28,
        3.2,
        0.28,
        size=10,
        color=GREEN,
        bold=True,
        font=MONO,
    )
    add_text(slide, title, 0.52, 0.61, 12.2, 0.55, size=29, bold=True)
    add_text(slide, subtitle, 0.52, 1.16, 12.1, 0.32, size=12.5, color=MUTED)
    add_text(
        slide,
        f"MARKHAND × CURSOR   {slide_number:02d}",
        10.65,
        7.12,
        2.15,
        0.2,
        size=8.5,
        color=MUTED,
        font=MONO,
        align=PP_ALIGN.RIGHT,
    )


def build_workflow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Slide 01 · Delivery workflow",
        "Từ ý tưởng đến sản phẩm với Cursor",
        "Human giữ quyền quyết định; Rule, Skill, Agent và CI đảm nhiệm các lớp khác nhau.",
        1,
    )

    cards = [
        ("01", "Brainstorm", "Mục tiêu • phương án", 0.52, 1.43, CYAN),
        ("02", "Plan", "Scope • rủi ro", 2.18, 1.30, PURPLE),
        ("03", "Phase", "Gate • dependency", 3.71, 1.37, AMBER),
        ("04", "Issue Ready", "AC • evidence", 5.31, 1.52, GREEN),
        ("05", "Orchestrate", "Phân vai độc lập", 7.06, 1.55, CYAN),
        ("06", "PR + CI", "Review • quality gate", 8.84, 1.52, PURPLE),
        ("07", "Done", "Evidence đầy đủ", 10.59, 1.42, GREEN),
    ]
    for number, title, subtitle, x, width, accent in cards:
        add_card(slide, number, title, subtitle, x, width, accent)

    for x in (2.00, 3.53, 5.13, 6.88, 8.66, 10.41):
        add_chevron(slide, x, 2.28)

    for x, width, accent in (
        (6.94, 1.02, CYAN),
        (8.02, 1.18, GREEN),
        (9.26, 1.02, PURPLE),
    ):
        add_rect(
            slide,
            x,
            3.52,
            width,
            0.48,
            fill=SURFACE_2,
            line=accent,
            radius=True,
        )
    for label, x, width in (
        ("Planner", 6.94, 1.02),
        ("Implementer", 8.02, 1.18),
        ("Reviewer", 9.26, 1.02),
    ):
        add_text(
            slide,
            label,
            x,
            3.52,
            width,
            0.48,
            size=10.5,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(7.84),
        Inches(3.20),
        Inches(7.84),
        Inches(3.38),
    )
    line.line.color.rgb = BORDER
    line.line.width = Pt(1.5)
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(7.45),
        Inches(3.38),
        Inches(9.77),
        Inches(3.38),
    )
    line.line.color.rgb = BORDER
    line.line.width = Pt(1.5)

    add_rect(slide, 0.52, 4.44, 12.29, 1.57, fill=SURFACE_2)
    legends = [
        (
            "RULE",
            "Guardrails",
            "Kiến trúc, bảo mật và giới hạn luôn đúng",
            0.82,
            GREEN,
        ),
        (
            "SKILL",
            "Workflow",
            "Cách thực hiện một loại công việc lặp lại",
            4.73,
            CYAN,
        ),
        (
            "HUMAN",
            "Approval",
            "Duyệt scope, plan, merge và release",
            8.64,
            AMBER,
        ),
    ]
    for label, title, body, x, accent in legends:
        add_text(
            slide,
            label,
            x,
            4.73,
            0.84,
            0.26,
            size=9,
            color=accent,
            bold=True,
            font=MONO,
        )
        add_text(slide, title, x, 5.02, 3.1, 0.27, size=14, bold=True)
        add_text(slide, body, x, 5.34, 3.18, 0.42, size=10.5, color=MUTED)

    add_text(
        slide,
        "Merge ≠ Done",
        0.55,
        6.35,
        2.05,
        0.4,
        size=17,
        color=GREEN,
        bold=True,
    )
    add_text(
        slide,
        "Done chỉ khi acceptance criteria và evidence đều hoàn tất.",
        2.34,
        6.35,
        7.6,
        0.4,
        size=14,
        color=WHITE,
    )


def add_repo_card(slide, x, accent, step, path, purpose):
    add_rect(slide, x, 1.85, 2.28, 1.36, fill=SURFACE)
    add_text(
        slide,
        step,
        x + 0.15,
        2.02,
        0.36,
        0.28,
        size=10,
        color=accent,
        bold=True,
        font=MONO,
    )
    add_text(
        slide,
        path,
        x + 0.15,
        2.33,
        1.98,
        0.36,
        size=12.5,
        bold=True,
        font=MONO,
        valign=MSO_ANCHOR.TOP,
    )
    add_text(
        slide,
        purpose,
        x + 0.15,
        2.74,
        1.98,
        0.28,
        size=10,
        color=MUTED,
        valign=MSO_ANCHOR.TOP,
    )


def build_repository_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(
        slide,
        "Slide 02 · Repository knowledge",
        "Repository là nguồn sự thật cho người và Agent",
        "Mỗi loại tài liệu có một nhiệm vụ: context → quyết định → delivery → automation.",
        2,
    )

    repo_cards = [
        (
            0.52,
            GREEN,
            "01 · CONTEXT",
            "AGENTS.md\nCLAUDE.md",
            "Kiến trúc • giới hạn",
        ),
        (
            3.02,
            CYAN,
            "02 · DECISIONS",
            "plans/reports/",
            "Brainstorm đã duyệt",
        ),
        (
            5.52,
            AMBER,
            "03 · ROADMAP",
            "plans/markhand-web/",
            "Phase • milestone • gate",
        ),
        (
            8.02,
            PURPLE,
            "04 · DELIVERY",
            "backlog/ + docs/",
            "Issue • Ready/Done",
        ),
        (
            10.52,
            GREEN,
            "05 · EXECUTION",
            ".cursor/ + .github/",
            "Rule • Skill • CI",
        ),
    ]
    for card in repo_cards:
        add_repo_card(slide, *card)
    for x in (2.78, 5.28, 7.78, 10.28):
        add_chevron(slide, x, 2.34)

    add_rect(slide, 0.52, 3.70, 12.28, 2.22, fill=SURFACE_2)
    add_text(
        slide,
        "Agent đọc",
        0.83,
        4.05,
        2.76,
        0.35,
        size=15,
        color=CYAN,
        bold=True,
    )
    add_text(
        slide,
        "Context + constraints\nIssue + acceptance criteria",
        0.83,
        4.45,
        2.76,
        0.72,
        size=12,
        color=WHITE,
        valign=MSO_ANCHOR.TOP,
    )
    add_chevron(slide, 3.72, 4.54, CYAN)

    add_text(
        slide,
        "Agent thực hiện",
        4.12,
        4.05,
        2.76,
        0.35,
        size=15,
        color=GREEN,
        bold=True,
    )
    add_text(
        slide,
        "Plan → code → review\nTheo Rule và Skill",
        4.12,
        4.45,
        2.76,
        0.72,
        size=12,
        color=WHITE,
        valign=MSO_ANCHOR.TOP,
    )
    add_chevron(slide, 7.01, 4.54, GREEN)

    add_text(
        slide,
        "Repository lưu lại",
        7.42,
        4.05,
        2.76,
        0.35,
        size=15,
        color=PURPLE,
        bold=True,
    )
    add_text(
        slide,
        "Commit + PR + CI\nevidence có version",
        7.42,
        4.45,
        2.76,
        0.72,
        size=12,
        color=WHITE,
        valign=MSO_ANCHOR.TOP,
    )
    add_chevron(slide, 10.31, 4.54, PURPLE)

    add_text(
        slide,
        "GitHub",
        10.72,
        4.05,
        1.58,
        0.35,
        size=15,
        color=AMBER,
        bold=True,
    )
    add_text(
        slide,
        "Milestone • Issue\nPR • Check",
        10.72,
        4.45,
        1.58,
        0.72,
        size=12,
        color=WHITE,
        valign=MSO_ANCHOR.TOP,
    )

    add_text(
        slide,
        "Agent không tự “nhớ” dự án.",
        0.55,
        6.35,
        3.30,
        0.4,
        size=17,
        color=GREEN,
        bold=True,
    )
    add_text(
        slide,
        "Repository cung cấp context, quy trình và bằng chứng có version control.",
        3.62,
        6.35,
        8.25,
        0.4,
        size=14,
        color=WHITE,
    )


def validate(prs):
    assert len(prs.slides) == 2
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            assert shape.left >= 0 and shape.top >= 0, (
                f"slide {slide_number}: negative shape position"
            )
            assert shape.left + shape.width <= prs.slide_width + Inches(0.01), (
                f"slide {slide_number}: shape exceeds slide width"
            )
            assert shape.top + shape.height <= prs.slide_height + Inches(0.01), (
                f"slide {slide_number}: shape exceeds slide height"
            )


def main():
    prs = Presentation()
    prs.slide_width = WIDE_WIDTH
    prs.slide_height = WIDE_HEIGHT
    prs.core_properties.title = "Cursor workflow — Markhand training"
    prs.core_properties.subject = "Two-slide workflow and repository knowledge overview"
    prs.core_properties.author = "Markhand"
    build_workflow_slide(prs)
    build_repository_slide(prs)
    validate(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()

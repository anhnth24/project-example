#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Tạo corpus ground-truth tiếng Việt cho đo độ chính xác.

Cùng MỘT đoạn text tham chiếu được nhúng vào docx/pptx/xlsx/csv/html để so
sánh độ chính xác trích xuất giữa các định dạng. Ảnh (chữ in) render từ text
để kiểm thử OCR Tesseract tiếng Việt.
"""
import os
import csv as csvmod

OUT = os.path.join(os.path.dirname(__file__), "vn_corpus")
os.makedirs(OUT, exist_ok=True)

# Đoạn tham chiếu giàu dấu tiếng Việt, số, dấu câu.
REF = (
    "Cộng hòa Xã hội Chủ nghĩa Việt Nam. "
    "Hôm nay là ngày 30 tháng 6 năm 2026. "
    "Tôi đang kiểm thử công cụ chuyển đổi tài liệu sang Markdown viết bằng Rust "
    "để đạt hiệu năng cao nhất. "
    "Tiếng Việt có nhiều dấu thanh: sắc, huyền, hỏi, ngã, nặng. "
    "Ví dụ các nguyên âm có dấu: ăn, ắt, ằng, ẳng, ẵng, ặng, ơ, ớ, ờ, ở, ỡ, ợ, ư, ứ, ừ, ử, ữ, ự. "
    "Giá trị hợp đồng là 1.234.567 đồng, tương đương khoảng năm mươi đô la Mỹ."
)
# Tách câu để đổ vào nhiều dòng/ô.
SENTENCES = [s.strip() + "." for s in REF.split(". ") if s.strip()]
SENTENCES[-1] = SENTENCES[-1].rstrip(".") + "."  # chuẩn lại dấu cuối

def w(name, data, mode="w"):
    p = os.path.join(OUT, name)
    with open(p, mode, encoding="utf-8") as f:
        f.write(data)
    print("  wrote", name)

# 0) ground truth
w("ref.txt", REF)

# 1) DOCX
from docx import Document
doc = Document()
for s in SENTENCES:
    doc.add_paragraph(s)
doc.save(os.path.join(OUT, "vn.docx"))
print("  wrote vn.docx")

# 2) PPTX (text box - kịch bản "pptx text")
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(5))
tf = tb.text_frame
tf.word_wrap = True
tf.text = SENTENCES[0]
for s in SENTENCES[1:]:
    tf.add_paragraph().text = s
prs.save(os.path.join(OUT, "vn.pptx"))
print("  wrote vn.pptx")

# 3) XLSX - mỗi câu một dòng, cột A
from openpyxl import Workbook
wb = Workbook()
ws = wb.active
for s in SENTENCES:
    ws.append([s])
wb.save(os.path.join(OUT, "vn.xlsx"))
print("  wrote vn.xlsx")

# 4) CSV - mỗi câu một dòng
with open(os.path.join(OUT, "vn.csv"), "w", encoding="utf-8", newline="") as f:
    wr = csvmod.writer(f)
    for s in SENTENCES:
        wr.writerow([s])
print("  wrote vn.csv")

# 5) HTML
html = "<html><head><meta charset='utf-8'><title>VN</title></head><body>"
for s in SENTENCES:
    html += f"<p>{s}</p>"
html += "</body></html>"
w("vn.html", html)

print("Xong corpus văn bản. Ảnh sẽ render bằng ImageMagick ở bước sau.")
